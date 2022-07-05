from curses.ascii import SP
import json
import os, sys
import argparse
import altair as alt
import sqlite3
import pandas as pd
from altair_saver import save
from pptx import Presentation
from pptx.util import Inches
from static.pd2ppt import df_to_table
from datetime import datetime
from datetime import timedelta
from PIL import Image

# ### Parsing arguments ################################# #

examples = '''examples:
  {prog}                    : Reports the data of the last completed month
  {prog} -y 2022            : Reports the data of the last completed month of 2022
  {prog} -m 2               : Reports the data of 2nd month of this year
  {prog} -day 22            : Reports the data of 22nd day of the last completed month
  {prog} -month 2 -d 22     : Reports the data of February 2nd this year
  {prog} -m 2 -year 2022    : Reports the data of February this year
  {prog} -d 2 -m 2 -y 2022  : Reports the data of 2nd February 2022
'''.format(prog="hvPptx.py")

def parse_args(args=sys.argv[1:]):
    parser = argparse.ArgumentParser(prog="hvPptx.py", description='Use the ' + "hvPptx.py" + ' command to create the report.', epilog=examples, formatter_class=argparse.RawDescriptionHelpFormatter, add_help=False)

    parser.add_argument('-help', '-h', action='help', default=argparse.SUPPRESS, help='Show this help message and exit.')
    parser.add_argument('-version', '-v', action='version', version='%(prog)s ' + hitachiConfig['version'], help='Show program\'s version number and exit.')

    group_main = parser.add_argument_group("usage")
    group_main.add_argument('-long', '-l', action='store_true', help='(Optional) Get long report.')
    group_main.add_argument('-day', '-d', type=int, metavar='<DAY>', help='(Optional) Specify the day of the data you want to report.')
    group_main.add_argument('-month', '-m', type=int, metavar='<MONTH>', help='(Optional) Specify the month of the data you want to report.')
    group_main.add_argument('-year', '-y', type=int, metavar='<YEAR>', help='(Optional) Specify the year of the data you want to report.')

    return parser, parser.parse_args(args)

def check_parse_args(parser):
    if args.year and not args.month: parser.error('-month should be provided when -year is used.')

def get_date():
    today = datetime.today()
    first = today.replace(day=1)
    if not len(sys.argv) > 1 or (args.long is not None and not len(sys.argv) > 2):      
        f_edate = first - timedelta(days=1)
        f_edate = f_edate.replace(hour=23).replace(minute=59).replace(second=59)
        f_sdate = f_edate.replace(day=1).replace(hour=00).replace(minute=00).replace(second=00)
    else:
        year = args.year if args.year is not None else datetime.today().strftime('%Y')
        if args.month is not None:
            month = args.month 
        else:
            t_date = first - timedelta(days=1)
            month = t_date.strftime('%m')
        if args.day is not None:
            day = args.day
            f_sdate = str(year) + str(month) + str(day) + ' 000000'
            f_sdate = datetime.strptime(f_sdate, '%Y%m%d %H%M%S')
            f_edate = f_sdate.replace(hour=23).replace(minute=59).replace(second=59)
        else:
            f_sdate = str(year) + str(month) + '01' + ' 000000'
            f_sdate = datetime.strptime(f_sdate, '%Y%m%d %H%M%S')
            next_month = f_sdate.replace(day=28) + timedelta(days=4)
            f_edate = next_month - timedelta(days=next_month.day)
            f_edate = f_edate.replace(hour=23).replace(minute=59).replace(second=59)
    return f_sdate, f_edate

def splitAdministratorData(columnsStr, columnsFloat):
    editedColStr = [i.split('.', -2)[0] for i in columnsStr]
    editedColFloat = []
    for item in columnsFloat:
        if 'InBytes' in item: item = item.split('InBytes')[0]
        if '.' in item: item = item.split('.')[-2]
        editedColFloat.append(item)
    return editedColStr, editedColFloat

def getDataFrame(source, fProc):
    source['date'] = pd.to_datetime(source['date'])
    mask = (source['date'] > sDate) & (source['date'] <= eDate)
    source = source.loc[mask]
    source = source[source[fProc['parameter']] != -1]
    maxdf = eval(fProc['melt']['max'])
    avgdf = eval(fProc['melt']['avg'])
    idVars = eval(fProc['melt']['max'][fProc['melt']['max'].find("(")+1:fProc['melt']['max'].find(")")])
    source = maxdf.merge(avgdf, how = 'inner', on = idVars)
    source = pd.melt(source, id_vars = idVars, value_vars = ['max', 'avg'])
    source['value'] = source['value'].round().astype(int)
    return source

def getChart(source, fProc):
    if fProc['column'] == 'port':
        singleWidth = 10
    else:
        singleWidth = 50
    base = alt.Chart(source).encode(
        x=alt.X('{}:O'.format(fProc['x']), axis=alt.Axis(labels=False, title= None, ticks=False)),
        y=alt.Y('{}:Q'.format(fProc['y']), axis=alt.Axis(title=fProc['ytitle'])),
        color='{}:N'.format(fProc['x']),
    )

    basePlot = alt.layer(
        base.mark_bar(),
        base.mark_text(dy=-10, align='center',).encode(text=fProc['y']),
    ).properties(
        width=singleWidth,
        height=150,
    )

    if len(source[fProc['column']].iloc[-1]) < 8:
        myplot = basePlot.facet(
            column='{}:N'.format(fProc['column']),
        )
    else:
        myplot = basePlot.facet(
            column=alt.Column('{}:N'.format(fProc['column']), header=alt.Header(labelAngle=90, labelAlign='left'))
        )

    return myplot

def getBar(source, hvDev, title):
    if '(GB)' in title:
        dx = 20
    else:
        dx = 10
    if hvDev:
        base = alt.Chart(source, title=title).encode(
            x='value:Q',
            y='metric:O',
            color='metric:N',
        )

        myplot = alt.layer(
            base.mark_bar(),
            base.mark_text(dx=dx).encode(text='value')
        ).properties(
            width=600,
            height=150,
        ).facet(
            row='{}:N'.format(hvDev),
        )
    else:
        base = alt.Chart(source, title=title).mark_bar().encode(
            x='value:Q',
            y="metric:O"
        )
        text = base.mark_text(
            align='left',
            baseline='middle',
            dx=3  
        ).encode(
            text='value:Q',
        )
        myplot = (base + text).properties(
            width=700,
            height=200,
        )
    return myplot

def getPlot(source, title, columnStr, columnFloat):
    if len(columnStr) > 2:
        colorize = columnStr[2]
    else:
        colorize = 'variable'
    myplot = alt.Chart(source, title=title).mark_line(interpolate='basis').encode(
        x = 'date:T',
        y = 'value:Q',
        color='{}:N'.format(colorize),
    ).properties(
        width=800,
        height=300,
        description=title
    )
    return myplot

def imgResize(file):
    image = Image.open(file)
    imgWidth, imgHeight = image.size
    if imgWidth <= 880 and imgHeight <= 400:
        canvasHeight = 400
        canvasWidth = 880
        newImage = Image.new(image.mode, (canvasWidth, canvasHeight), 255)
        newImage.paste(image, (0, (canvasHeight - imgHeight) // 2))
    else:
        if imgWidth/imgHeight > 2.2:
            canvasHeight = int(imgWidth / 2.2)
            canvasWidth = imgWidth
        else:
            canvasWidth = int(imgHeight * 2.2)
            canvasHeight = imgHeight
        newImage = Image.new(image.mode, (canvasWidth, canvasHeight), 255)
        newImage.paste(image, ((canvasWidth - imgWidth) // 2, (canvasHeight - imgHeight) // 2))
    newImage.save(fp=file)

# ### MAIN Part ######################################### #

if __name__ == "__main__":
    try:
        with open('conf/hvData.json') as dataFile:
            hitachiData = json.load(dataFile)
        dataFile.close()
        with open('conf/hvConf.json') as confFile:
            hitachiConfig = json.load(confFile)
        confFile.close()
        with open('conf/hvProc.json') as procFile:
            hitachiProcess = json.load(procFile)
        confFile.close()
    except:
        print('!!! hvData.json or hvConf.json file cannot be opened!')
        raise SystemExit()

    parser, args = parse_args()
    check_parse_args(parser)

    alt.renderers.enable('default')
    alt.data_transformers.disable_max_rows()
    sDate, eDate = get_date()
    dbFile = 'data/db/hv.db'
    prs = Presentation('static/hvTemplate.pptx')
    top = Inches(1)
    left = Inches(0.25)
    width = Inches(9.5)
    height = Inches(0.3)
    strsDate = sDate.strftime("%Y-%m-%d")
    streDate = eDate.strftime("%Y-%m-%d")
    if strsDate == streDate:
        titleDate = strsDate
    else:
        titleDate = strsDate + ' - ' + streDate

    if os.path.isfile(dbFile):
        conn = sqlite3.connect(dbFile)
    else:
        print('!!! An error occured while opening the database ({})!'.format(dbFile))

    df = pd.read_sql_query('SELECT * FROM hvStorages1', conn)
    storages = df.sort_values(by=['storageSystemId'])['storageSystemId'].drop_duplicates().to_list()

    for sData in hitachiData:
        for sTable in sData['data']:
            dbTable = sData['table'] + sTable['id']
            plotdf = pd.DataFrame()
            if sData['type'] == 'administrator' and sTable['type'] == 'table':
                for storageSystemId in storages:
                    stdf = pd.read_sql_query('SELECT * FROM "{}" where storageSystemId="{}" ORDER by date'.format(dbTable, storageSystemId), conn)
                    if not stdf.empty:
                        lastDate = stdf['date'].iloc[-1]
                        stdf = stdf[stdf.date.isin([lastDate])].reset_index(drop=True)
                        if plotdf.empty:
                            plotdf = stdf
                        else:
                            plotdf = pd.concat([plotdf, stdf], ignore_index=True)
                if not plotdf.empty:
                    slide = prs.slides.add_slide(prs.slide_layouts[2])
                    title = slide.shapes.title
                    title.text = sTable['title']
                    plotdf.drop(columns=['date'], inplace=True)
                    if dbTable == 'hvParityGroups1':
                        plotdf = plotdf.groupby(['storageSystemId','raidLevel','raidLayout','encryption','compression','diskSpec']).size().reset_index(name='Count')
                    df_to_table(slide, plotdf, left, top, width, height)
    
    for sProc in hitachiProcess:
        chart = ''
        if sProc['type'] != 'series':
            plotdf = pd.read_sql_query(sProc['query'], conn)
            if sProc['type'] == 'allofseries':
                plotdf = getDataFrame(plotdf, sProc)
            else:
                valueVars = []
                for k,v in sProc['melt'].items():
                    if v:
                        plotdf[k] = eval(v)
                        plotdf[k] = plotdf[k].round().astype(int)
                    valueVars.append(k)
                if len(sProc['melt']) > 1:
                    plotdf = pd.melt(plotdf, id_vars = ['storageSystemId'], value_vars = valueVars)
            
            if not plotdf.empty:
                chart = getChart(plotdf, sProc)

            if chart:
                save(chart, 'data/tmp/{}.png'.format(sProc['filename']))
                imgResize('data/tmp/{}.png'.format(sProc['filename']))
                slide = prs.slides.add_slide(prs.slide_layouts[3])
                title = slide.shapes.title
                title.text = sProc['title']
                placeholder = slide.placeholders[1]
                placeholder.insert_picture('data/tmp/{}.png'.format(sProc['filename']))
        else:
            for storageSystemId in storages:
                chart = ''
                plotdf = pd.read_sql_query(sProc['query'].format(storageSystemId), conn)
                plotdf = getDataFrame(plotdf, sProc)
                if not plotdf.empty:
                    chart = getChart(plotdf, sProc)
                if chart:
                    save(chart, 'data/tmp/{}-{}.png'.format(sProc['filename'], storageSystemId))
                    imgResize('data/tmp/{}-{}.png'.format(sProc['filename'], storageSystemId))
                    slide = prs.slides.add_slide(prs.slide_layouts[3])
                    title = slide.shapes.title
                    title.text = storageSystemId + ': ' + sProc['title']
                    placeholder = slide.placeholders[1]
                    placeholder.insert_picture('data/tmp/{}-{}.png'.format(sProc['filename'], storageSystemId))

    if args.long is True:
        for sData in hitachiData:
            for sTable in sData['data']:
                dbTable = sData['table'] + sTable['id']
                plotdf = pd.DataFrame()
                if sData['type'] == 'analyzer' or sTable['type'] == 'plot':
                    for storageSystemId in storages:
                        plotdf = pd.read_sql_query('SELECT * FROM "{}" where storageSystemId="{}" ORDER by date'.format(dbTable, storageSystemId), conn)
                        chart = ''
                        if sData['type'] == 'analyzer':
                            if sTable['type'] != 'none':
                                if sTable['type'] == 'monthly' or (sTable['type'] == 'daily' and args.day is not None):
                                    plotdf['date'] = pd.to_datetime(plotdf['date'])
                                    plotdf = pd.melt(plotdf, id_vars =sTable['parameter']['columnsStr'], value_vars = sTable['parameter']['columnsFloat'])
                        
                                    mask = (plotdf['date'] > sDate) & (plotdf['date'] <= eDate)
                                    plotdf = plotdf.loc[mask]
                                    chart = getPlot(plotdf, sTable['title'], sTable['parameter']['columnsStr'], sTable['parameter']['columnsFloat'])
                                    totalTitle = storageSystemId + ': ' + sTable['title'] + ': ' + titleDate
                        else:
                            adminColStr, adminColFloat = splitAdministratorData(sTable['parameter']['columnsStr'], sTable['parameter']['columnsFloat'])
                            if '(GB)' in sTable['title']:
                                plotdf[adminColFloat] = plotdf[adminColFloat].astype(int)
                            plotdf = pd.melt(plotdf, id_vars = adminColStr, value_vars = adminColFloat)
                            plotdf = plotdf.rename(columns = {'variable':'metric'})
                            if len(adminColStr) > 2:
                                chart = getBar(plotdf, adminColStr[2], sTable['title'])
                            else:
                                chart = getBar(plotdf, None, sTable['title'])
                            totalTitle = storageSystemId + ': ' + sTable['title']

                        if chart:
                            save(chart, 'data/tmp/{}-{}.png'.format(dbTable, storageSystemId))
                            imgResize('data/tmp/{}-{}.png'.format(dbTable, storageSystemId))
                            slide = prs.slides.add_slide(prs.slide_layouts[3])
                            title = slide.shapes.title
                            title.text = totalTitle
                            placeholder = slide.placeholders[1]
                            placeholder.insert_picture('data/tmp/{}-{}.png'.format(dbTable, storageSystemId))
                        
    prs.slides.add_slide(prs.slide_layouts[12])
    prs.save('data/tmp/hvReport.pptx')
    for f in os.listdir('data/tmp'):
        if f.split('.')[-1] == 'png':
            os.remove('data/tmp/{}'.format(f))

    